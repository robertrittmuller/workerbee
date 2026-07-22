import asyncio
import json
import logging
import os
import shutil
import uuid
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import joinedload

from app.config import settings
from app.models import Execution, ExecutionLog, Artifact, File as FileModel
from app.work_packs import format_work_pack_instructions
from app.opencode_client import opencode_client

logger = logging.getLogger(__name__)

TOOL_FREE_MODE = {
    tool_name: False
    for tool_name in (
        "bash",
        "write",
        "edit",
        "apply_patch",
        "read",
        "glob",
        "grep",
        "list",
        "task",
        "webfetch",
        "websearch",
        "todowrite",
        "todoread",
        "question",
        "skill",
        "lsp",
    )
}
MAX_INLINE_SOURCE_CHARS = 60_000


def _extract_source_text(path: Path) -> str:
    """Extract readable content from common business file formats for fallback mode."""
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".xml", ".html", ".htm", ".yaml", ".yml"}:
            return path.read_text(encoding="utf-8", errors="replace")

        if suffix == ".pdf":
            reader = PdfReader(path)
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)

        if suffix == ".docx":
            document = Document(path)
            blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
            for table in document.tables:
                for row in table.rows:
                    blocks.append("\t".join(cell.text for cell in row.cells))
            return "\n".join(blocks)

        if suffix == ".xlsx":
            workbook = load_workbook(path, read_only=True, data_only=True)
            blocks: list[str] = []
            for worksheet in workbook.worksheets:
                blocks.append(f"[Sheet: {worksheet.title}]")
                for row in worksheet.iter_rows(values_only=True):
                    blocks.append("\t".join("" if value is None else str(value) for value in row))
                    if sum(len(item) for item in blocks) >= MAX_INLINE_SOURCE_CHARS:
                        break
            workbook.close()
            return "\n".join(blocks)

        if suffix == ".pptx":
            presentation = Presentation(path)
            blocks = []
            for slide_number, slide in enumerate(presentation.slides, start=1):
                blocks.append(f"[Slide {slide_number}]")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if isinstance(text, str) and text.strip():
                        blocks.append(text)
            return "\n".join(blocks)
    except Exception as exc:
        logger.warning("Could not extract fallback text from %s: %s", path, exc)
    return ""


def _opencode_error(result: Any) -> str | None:
    if not isinstance(result, dict):
        return None
    error = result.get("info", {}).get("error")
    if not error:
        return None
    if isinstance(error, dict):
        data = error.get("data")
        if isinstance(data, dict) and isinstance(data.get("message"), str):
            return data["message"]
        if isinstance(error.get("message"), str):
            return error["message"]
    return str(error)

async def execute_agent(
    execution_id: Any,
    agent_config: dict[str, Any],
    task_prompt: str,
    input_files: list[dict[str, str]],
    output_config: dict[str, Any],
    opencode_agent: str = "general",
) -> dict[str, Any]:
    from app.database import async_session_maker
    async with async_session_maker() as db:
        res = await run_agent(execution_id, db, task_prompt=task_prompt, opencode_agent=opencode_agent, input_files=input_files)
    
    if not res:
        return {"success": False, "error": "Agent failed to return a valid response."}
    return res

async def run_agent(
    execution_id: Any,
    db: AsyncSession,
    task_prompt: str = "Do your task.",
    opencode_agent: str = "general",
    input_files: list[dict[str, str]] | None = None
):
    execution_uuid = execution_id if isinstance(execution_id, uuid.UUID) else uuid.UUID(str(execution_id))

    # Retrieve execution
    stmt = select(Execution).options(joinedload(Execution.agent)).where(Execution.id == execution_uuid)
    result = await db.execute(stmt)
    execution = result.scalar_one_or_none()
    
    if not execution:
        logger.error(f"Execution {execution_id} not found")
        return
        
    try:
        execution.status = "running"
        await db.commit()

        # Step 1: Create OpenCode session
        session_data = await opencode_client.create_session(title=f"WorkerBee {execution_id}")
        session_id = session_data["id"]
        
        execution.opencode_session_id = session_id
        await db.commit()
        
        await log_msg(db, execution_uuid, "System", "OpenCode session created.")

        # Step 2: Prepare workspace
        workspace_dir = Path(settings.opencode_workspace_root) / f"executions/{execution_id}"
        workspace_dir.mkdir(parents=True, exist_ok=True)
        output_dir = workspace_dir / "output"
        output_dir.mkdir(parents=True, exist_ok=True)
        inline_sources: list[str] = []
        
        if input_files:
            for file_info in input_files:
                src_path = Path(file_info["storage_path"])
                if not src_path.is_absolute():
                    src_path = Path.cwd() / src_path
                if src_path.exists():
                    dest_path = workspace_dir / file_info["filename"]
                    shutil.copy2(src_path, dest_path)
                    extracted_text = _extract_source_text(dest_path).strip()
                    if extracted_text:
                        remaining_chars = MAX_INLINE_SOURCE_CHARS - sum(
                            len(source) for source in inline_sources
                        )
                        if remaining_chars > 0:
                            inline_sources.append(
                                f"SOURCE: {file_info['filename']}\n"
                                f"{extracted_text[:remaining_chars]}"
                            )
                    await log_msg(db, execution_uuid, "System", f"Copied input file: {file_info['filename']}")
                else:
                    await log_msg(db, execution_uuid, "System", f"Warning: Input file not found at {src_path}")
        
        # Step 3: Send prompt
        agent_config = execution.agent.config
        template_id = agent_config.get("template", {}).get("id", "unknown")
        template_markdown = agent_config.get("template", {}).get("markdown_files", [])
        template_instructions = "\n\n".join(
            item.get("content", "").strip()
            for item in template_markdown
            if isinstance(item, dict) and isinstance(item.get("content"), str)
        ).strip()
        work_pack_instructions = format_work_pack_instructions(agent_config.get("work_pack"))
        
        prompt = (
            f"WORKERBEE TASK INSTRUCTIONS\n{template_instructions}\n\n"
            f"{work_pack_instructions}\n\n"
            f"USER REQUEST\n{task_prompt}\n\n"
            f"WORKSPACE\nYour private workspace is: {workspace_dir}\n"
            f"Input files are located directly inside that directory.\n"
            f"Write every requested deliverable to this exact output directory: {output_dir}\n"
            "Do not leave final deliverables elsewhere in the workspace. Create useful, polished files "
            "with descriptive filenames. In your final response, briefly summarize what you created."
        )
        
        # Send to OpenCode using a supported built-in agent (OpenCode only has: build, explore, general, plan)
        await log_msg(db, execution_uuid, "System", f"Starting '{opencode_agent}' agent execution (Template: {template_id}) ...")
        
        # This will block until the model finishes its task
        result = await opencode_client.send_prompt(session_id, prompt, agent_name=opencode_agent)
        tool_free_fallback = False
        engine_error = _opencode_error(result)
        if engine_error and "tool" in engine_error.lower():
            tool_free_fallback = True
            inline_source_text = (
                "\n\n".join(inline_sources)
                if inline_sources
                else "[No readable source text was extracted.]"
            )
            await log_msg(
                db,
                execution_uuid,
                "System",
                "The selected model cannot use workspace tools; retrying with a direct deliverable.",
            )
            fallback_prompt = (
                f"{prompt}\n\n"
                "FALLBACK MODE: No tools are available. Do not claim to create or edit files. "
                "Return the complete requested deliverable directly in your final response, using the "
                "requested format and including all final content.\n\n"
                "SOURCE MATERIAL EXTRACTED FROM THE USER'S FILES:\n"
                f"{inline_source_text}"
            )
            result = await opencode_client.send_prompt(
                session_id,
                fallback_prompt,
                agent_name=opencode_agent,
                tools=TOOL_FREE_MODE,
            )
            engine_error = _opencode_error(result)
        
        # Step 4: Harvest artifacts from workspace_dir
        # Assuming artifacts would be discovered and registered here...
        
        await db.commit()
        await log_msg(db, execution_uuid, "System", "Execution completed successfully.")
        
        # Extract response text
        if engine_error:
            return {
                "success": False,
                "error": engine_error,
                "opencode_result": result,
            }

        content = "No output text provided by agent."
        if isinstance(result, dict):
            # Check OpenCode return format for text
            parts = result.get("parts", [])
            for part in reversed(parts):
                if part.get("type") == "text":
                    content = part.get("text", content)
                    break
        
        return {
            "success": True,
            "messages": [{"role": "assistant", "content": content}],
            "opencode_result": result,
            "tool_free_fallback": tool_free_fallback,
        }
        
    except Exception as e:
        logger.error(f"Execution {execution_id} failed: {e}", exc_info=True)
        await log_msg(db, execution_uuid, "System", f"Execution failed: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }

async def log_msg(db: AsyncSession, exec_id: Any, role: str, msg: str):
    execution_uuid = exec_id if isinstance(exec_id, uuid.UUID) else uuid.UUID(str(exec_id))
    logger.info(f"[{exec_id}] {role}: {msg}")
    db.add(ExecutionLog(
        execution_id=execution_uuid,
        level="info",
        message=f"{role}: {msg}"
    ))
    await db.commit()
    
