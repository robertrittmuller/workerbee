"""Bounded, read-only previews for common business source files."""

from __future__ import annotations

import csv
import json
import zipfile
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.schemas import FilePreviewResponse, FilePreviewTable

MAX_PREVIEW_FILE_BYTES = 20 * 1024 * 1024
MAX_TEXT_CHARS = 60_000
MAX_TABLE_ROWS = 80
MAX_TABLE_COLUMNS = 30
MAX_TABLES = 8
MAX_PDF_PAGES = 20
MAX_SLIDES = 30
MAX_OFFICE_UNCOMPRESSED_BYTES = 120 * 1024 * 1024

TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".html",
    ".htm",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".webp"}


def _unavailable(detail: str) -> FilePreviewResponse:
    return FilePreviewResponse(kind="unavailable", detail=detail)


def _bounded_text(value: str) -> tuple[str, bool]:
    if len(value) <= MAX_TEXT_CHARS:
        return value, False
    return value[:MAX_TEXT_CHARS].rstrip(), True


def _read_bounded_text(path: Path) -> tuple[str, bool]:
    with path.open("rb") as handle:
        content = handle.read(MAX_TEXT_CHARS + 1)
    decoded = content.decode("utf-8", errors="replace")
    return _bounded_text(decoded)


def _format_cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    return str(value)


def _tabular_preview(name: str, rows: list[list[str]], truncated: bool) -> FilePreviewTable:
    normalized = [row[:MAX_TABLE_COLUMNS] for row in rows]
    return FilePreviewTable(name=name, rows=normalized, truncated=truncated)


def _preview_delimited(path: Path, delimiter: str) -> FilePreviewResponse:
    rows: list[list[str]] = []
    truncated = False
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row_number, row in enumerate(reader):
            if row_number >= MAX_TABLE_ROWS:
                truncated = True
                break
            if len(row) > MAX_TABLE_COLUMNS:
                truncated = True
            rows.append([str(value) for value in row[:MAX_TABLE_COLUMNS]])
    return FilePreviewResponse(
        kind="table",
        tables=[_tabular_preview(path.stem or "Data", rows, truncated)],
        truncated=truncated,
        detail="A bounded table preview. The source file is unchanged.",
    )


def _office_archive_is_bounded(path: Path) -> bool:
    try:
        with zipfile.ZipFile(path) as archive:
            total = 0
            for item in archive.infolist():
                total += item.file_size
                if total > MAX_OFFICE_UNCOMPRESSED_BYTES:
                    return False
        return True
    except zipfile.BadZipFile:
        return False


def _preview_docx(path: Path) -> FilePreviewResponse:
    if not _office_archive_is_bounded(path):
        return _unavailable("This document is damaged or expands beyond the safe preview limit.")
    document = Document(path)
    blocks: list[str] = []
    truncated = False
    current_chars = 0

    def append_block(value: str) -> bool:
        nonlocal current_chars, truncated
        clean = value.strip()
        if not clean:
            return True
        remaining = MAX_TEXT_CHARS - current_chars
        if remaining <= 0:
            truncated = True
            return False
        if len(clean) > remaining:
            blocks.append(clean[:remaining].rstrip())
            current_chars = MAX_TEXT_CHARS
            truncated = True
            return False
        blocks.append(clean)
        current_chars += len(clean) + 2
        return True

    for paragraph in document.paragraphs:
        if not append_block(paragraph.text):
            break
    if not truncated:
        for table in document.tables:
            if not append_block("\n".join("\t".join(cell.text for cell in row.cells) for row in table.rows)):
                break
    return FilePreviewResponse(
        kind="text",
        text="\n\n".join(blocks),
        truncated=truncated,
        detail="Readable document text. Layout, comments, and tracked changes may not appear here.",
    )


def _preview_xlsx(path: Path) -> FilePreviewResponse:
    if not _office_archive_is_bounded(path):
        return _unavailable("This workbook is damaged or expands beyond the safe preview limit.")
    workbook = load_workbook(path, read_only=True, data_only=True)
    tables: list[FilePreviewTable] = []
    overall_truncated = len(workbook.worksheets) > MAX_TABLES
    try:
        for worksheet in workbook.worksheets[:MAX_TABLES]:
            rows: list[list[str]] = []
            sheet_truncated = False
            for row_number, values in enumerate(worksheet.iter_rows(values_only=True)):
                if row_number >= MAX_TABLE_ROWS:
                    sheet_truncated = True
                    break
                row_values = [_format_cell(value) for value in values]
                if len(row_values) > MAX_TABLE_COLUMNS:
                    sheet_truncated = True
                rows.append(row_values[:MAX_TABLE_COLUMNS])
            tables.append(_tabular_preview(worksheet.title, rows, sheet_truncated))
            overall_truncated = overall_truncated or sheet_truncated
    finally:
        workbook.close()
    return FilePreviewResponse(
        kind="table",
        tables=tables,
        truncated=overall_truncated,
        detail="Calculated cell values from a bounded workbook preview. Formulas and formatting are not changed.",
    )


def _preview_pptx(path: Path) -> FilePreviewResponse:
    if not _office_archive_is_bounded(path):
        return _unavailable("This presentation is damaged or expands beyond the safe preview limit.")
    presentation = Presentation(path)
    blocks: list[str] = []
    truncated = len(presentation.slides) > MAX_SLIDES
    current_chars = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number > MAX_SLIDES:
            break
        slide_text = [
            text.strip()
            for shape in slide.shapes
            if isinstance((text := getattr(shape, "text", "")), str) and text.strip()
        ]
        block = f"Slide {slide_number}\n" + ("\n".join(slide_text) or "No readable text")
        remaining = MAX_TEXT_CHARS - current_chars
        if remaining <= 0:
            truncated = True
            break
        if len(block) > remaining:
            blocks.append(block[:remaining].rstrip())
            truncated = True
            break
        blocks.append(block)
        current_chars += len(block) + 2
    return FilePreviewResponse(
        kind="text",
        text="\n\n".join(blocks),
        truncated=truncated,
        detail="Readable slide text. Visual layout, charts, notes, and animations may not appear here.",
    )


def _preview_pdf(path: Path) -> FilePreviewResponse:
    reader = PdfReader(path)
    page_count = len(reader.pages)
    blocks: list[str] = []
    truncated = page_count > MAX_PDF_PAGES
    current_chars = 0
    for page_number, page in enumerate(reader.pages, start=1):
        if page_number > MAX_PDF_PAGES:
            break
        text = (page.extract_text() or "").strip()
        block = f"Page {page_number}\n" + (text or "No readable text on this page")
        remaining = MAX_TEXT_CHARS - current_chars
        if remaining <= 0:
            truncated = True
            break
        if len(block) > remaining:
            blocks.append(block[:remaining].rstrip())
            truncated = True
            break
        blocks.append(block)
        current_chars += len(block) + 2
    return FilePreviewResponse(
        kind="text",
        text="\n\n".join(blocks),
        page_count=page_count,
        truncated=truncated,
        detail="Extracted PDF text. Scanned pages, complex layout, and annotations may not appear here.",
    )


def build_file_preview(
    path: Path,
    *,
    filename: str,
    content_type: str,
    file_size: int,
) -> FilePreviewResponse:
    """Build a bounded preview without mutating or sending the source file."""
    if file_size > MAX_PREVIEW_FILE_BYTES:
        return _unavailable("This file is too large to preview safely. You can still use it in a task or save a copy.")
    if not path.exists() or not path.is_file():
        return _unavailable("This source file is unavailable.")

    suffix = Path(filename).suffix.lower()
    try:
        if suffix in IMAGE_SUFFIXES or content_type.startswith("image/"):
            return FilePreviewResponse(
                kind="image",
                detail="Image preview. The original file remains unchanged.",
            )
        if suffix in {".csv", ".tsv"}:
            return _preview_delimited(path, "\t" if suffix == ".tsv" else ",")
        if suffix in TEXT_SUFFIXES or content_type.startswith("text/"):
            text, truncated = _read_bounded_text(path)
            if suffix == ".json":
                try:
                    text = json.dumps(json.loads(text), indent=2, ensure_ascii=False)
                    text, formatted_truncated = _bounded_text(text)
                    truncated = truncated or formatted_truncated
                except json.JSONDecodeError:
                    pass
            return FilePreviewResponse(
                kind="text",
                text=text,
                truncated=truncated,
                detail="Plain-text preview. The source file is unchanged.",
            )
        if suffix == ".pdf":
            return _preview_pdf(path)
        if suffix == ".docx":
            return _preview_docx(path)
        if suffix == ".xlsx":
            return _preview_xlsx(path)
        if suffix == ".pptx":
            return _preview_pptx(path)
        if suffix in {".doc", ".xls", ".ppt"}:
            return _unavailable("Preview is not available for this older Office format. You can still use it in a task or save a copy.")
    except Exception:
        return _unavailable("WorkerBee could not read this file safely. The original file has not been changed.")

    return _unavailable("Preview is not available for this file type. You can still use it in a task or save a copy.")
