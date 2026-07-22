"""Deterministic PowerPoint rendering for WorkerBee presentation work packs."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


class PresentationRenderError(ValueError):
    """Raised when a deck specification cannot be safely rendered."""


THEMES = {
    "Executive dark": {
        "background": "171717",
        "surface": "292524",
        "primary": "FAFAF9",
        "secondary": "D6D3D1",
        "accent": "F59E0B",
        "muted": "78716C",
    },
    "Warm editorial": {
        "background": "F7F3EA",
        "surface": "FFFFFF",
        "primary": "292524",
        "secondary": "57534E",
        "accent": "B45309",
        "muted": "A8A29E",
    },
    "Clean light": {
        "background": "F8FAFC",
        "surface": "FFFFFF",
        "primary": "0F172A",
        "secondary": "475569",
        "accent": "0369A1",
        "muted": "94A3B8",
    },
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _clean_text(value: Any, *, field: str, maximum: int = 600) -> str:
    if not isinstance(value, str):
        raise PresentationRenderError(f"{field} must be text.")
    normalized = " ".join(value.replace("\x00", "").split()).strip()
    if not normalized:
        raise PresentationRenderError(f"{field} is required.")
    return normalized[:maximum]


def _optional_text(value: Any, maximum: int = 600) -> str:
    if not isinstance(value, str):
        return ""
    return " ".join(value.replace("\x00", "").split()).strip()[:maximum]


def _text_list(value: Any, *, maximum_items: int = 8, maximum_length: int = 280) -> list[str]:
    if not isinstance(value, list):
        return []
    items: list[str] = []
    for raw_item in value[:maximum_items]:
        item = _optional_text(raw_item, maximum_length)
        if item:
            items.append(item)
    return items


def _add_textbox(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    color: str,
    bold: bool = False,
    font: str = "Aptos",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    margin: float = 0,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)
    return shape


def _add_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_slide_chrome(
    slide: Any,
    *,
    theme: dict[str, str],
    deck_title: str,
    slide_number: int,
    slide_count: int,
    sources: list[str],
) -> None:
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(0.42),
        Inches(0.08),
        Inches(0.34),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme["accent"])
    accent.line.fill.background()
    _add_textbox(
        slide,
        left=0.78,
        top=0.39,
        width=7.5,
        height=0.35,
        text=deck_title.upper(),
        size=9,
        color=theme["muted"],
        bold=True,
    )
    _add_textbox(
        slide,
        left=11.75,
        top=0.39,
        width=0.95,
        height=0.35,
        text=f"{slide_number:02d} / {slide_count:02d}",
        size=9,
        color=theme["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    if sources:
        _add_textbox(
            slide,
            left=0.65,
            top=7.05,
            width=11.95,
            height=0.25,
            text="Sources: " + " · ".join(sources[:3]),
            size=7,
            color=theme["muted"],
        )


def _add_title_slide(slide: Any, slide_data: dict[str, Any], theme: dict[str, str]) -> None:
    title = _clean_text(slide_data.get("title"), field="title slide title", maximum=180)
    subtitle = _optional_text(slide_data.get("subtitle"), 300)
    eyebrow = _optional_text(slide_data.get("eyebrow"), 80) or "WORKERBEE BRIEFING"
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(0.85),
        Inches(0.12),
        Inches(5.75),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme["accent"])
    accent.line.fill.background()
    _add_textbox(
        slide,
        left=1.15,
        top=1.2,
        width=10.8,
        height=0.4,
        text=eyebrow.upper(),
        size=11,
        color=theme["accent"],
        bold=True,
    )
    _add_textbox(
        slide,
        left=1.15,
        top=2.0,
        width=10.8,
        height=2.65,
        text=title,
        size=50,
        color=theme["primary"],
        bold=True,
    )
    if subtitle:
        _add_textbox(
            slide,
            left=1.15,
            top=4.85,
            width=9.8,
            height=1.1,
            text=subtitle,
            size=24,
            color=theme["secondary"],
        )


def _add_section_slide(slide: Any, slide_data: dict[str, Any], theme: dict[str, str]) -> None:
    title = _clean_text(slide_data.get("title"), field="section title", maximum=160)
    subtitle = _optional_text(slide_data.get("subtitle"), 260)
    _add_textbox(
        slide,
        left=0.75,
        top=1.45,
        width=2.1,
        height=1.5,
        text=_optional_text(slide_data.get("section_number"), 12) or "—",
        size=52,
        color=theme["accent"],
        bold=True,
    )
    _add_textbox(
        slide,
        left=2.65,
        top=1.55,
        width=9.8,
        height=1.7,
        text=title,
        size=42,
        color=theme["primary"],
        bold=True,
    )
    if subtitle:
        _add_textbox(
            slide,
            left=2.68,
            top=3.45,
            width=8.8,
            height=1.2,
            text=subtitle,
            size=24,
            color=theme["secondary"],
        )


def _add_bullets(frame: Any, bullets: list[str], theme: dict[str, str], *, size: int = 18) -> None:
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.12)
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = _rgb(theme["secondary"])
        paragraph.space_after = Pt(12)
        paragraph.text = f"•  {bullet}"


def _add_content_slide(slide: Any, slide_data: dict[str, Any], theme: dict[str, str]) -> None:
    title = _clean_text(slide_data.get("title"), field="content slide title", maximum=160)
    takeaway = _optional_text(slide_data.get("takeaway"), 260)
    bullets = _text_list(slide_data.get("bullets"), maximum_items=6)
    _add_textbox(
        slide,
        left=0.75,
        top=0.95,
        width=11.7,
        height=1.05,
        text=title,
        size=35,
        color=theme["primary"],
        bold=True,
    )
    if takeaway:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.75),
            Inches(2.15),
            Inches(11.7),
            Inches(1.05),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(theme["surface"])
        box.line.color.rgb = _rgb(theme["accent"])
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = Inches(0.28)
        frame.margin_right = Inches(0.28)
        paragraph = frame.paragraphs[0]
        paragraph.text = takeaway
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = _rgb(theme["secondary"])
        paragraph.alignment = PP_ALIGN.CENTER
    body = slide.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.2), Inches(2.95))
    _add_bullets(body.text_frame, bullets or ["No supported details were supplied."], theme)


def _add_metrics_slide(slide: Any, slide_data: dict[str, Any], theme: dict[str, str]) -> None:
    title = _clean_text(slide_data.get("title"), field="metrics slide title", maximum=160)
    metrics = slide_data.get("metrics")
    if not isinstance(metrics, list) or not metrics:
        raise PresentationRenderError("A metrics slide requires metrics.")
    _add_textbox(
        slide,
        left=0.75,
        top=0.95,
        width=11.7,
        height=1.05,
        text=title,
        size=35,
        color=theme["primary"],
        bold=True,
    )
    normalized: list[tuple[str, str, str]] = []
    for index, metric in enumerate(metrics[:4]):
        if not isinstance(metric, dict):
            continue
        value = _clean_text(metric.get("value"), field=f"metric {index + 1} value", maximum=40)
        label = _clean_text(metric.get("label"), field=f"metric {index + 1} label", maximum=80)
        context = _optional_text(metric.get("context"), 120)
        normalized.append((value, label, context))
    if not normalized:
        raise PresentationRenderError("A metrics slide requires valid metrics.")
    card_width = (11.7 - (len(normalized) - 1) * 0.22) / len(normalized)
    for index, (value, label, context) in enumerate(normalized):
        left = 0.75 + index * (card_width + 0.22)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(2.15),
            Inches(card_width),
            Inches(3.55),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(theme["surface"])
        card.line.color.rgb = _rgb(theme["muted"])
        _add_textbox(slide, left=left + 0.25, top=2.55, width=card_width - 0.5, height=0.9, text=value, size=32, color=theme["accent"], bold=True)
        _add_textbox(slide, left=left + 0.25, top=3.55, width=card_width - 0.5, height=0.7, text=label, size=18, color=theme["primary"], bold=True)
        if context:
            _add_textbox(slide, left=left + 0.25, top=4.45, width=card_width - 0.5, height=0.9, text=context, size=16, color=theme["secondary"])


def _add_comparison_slide(slide: Any, slide_data: dict[str, Any], theme: dict[str, str]) -> None:
    title = _clean_text(slide_data.get("title"), field="comparison slide title", maximum=160)
    columns = slide_data.get("columns")
    if not isinstance(columns, list) or len(columns) < 2:
        raise PresentationRenderError("A comparison slide requires at least two columns.")
    _add_textbox(slide, left=0.75, top=0.95, width=11.7, height=1.05, text=title, size=35, color=theme["primary"], bold=True)
    normalized = [column for column in columns[:3] if isinstance(column, dict)]
    card_width = (11.7 - (len(normalized) - 1) * 0.25) / len(normalized)
    for index, column in enumerate(normalized):
        heading = _clean_text(column.get("heading"), field=f"comparison column {index + 1}", maximum=80)
        bullets = _text_list(column.get("bullets"), maximum_items=5, maximum_length=160)
        left = 0.75 + index * (card_width + 0.25)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.55), Inches(card_width), Inches(4.1))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(theme["surface"])
        card.line.color.rgb = _rgb(theme["muted"])
        _add_textbox(slide, left=left + 0.25, top=2.82, width=card_width - 0.5, height=0.7, text=heading, size=20, color=theme["primary"], bold=True)
        body = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.62), Inches(card_width - 0.4), Inches(2.7))
        _add_bullets(body.text_frame, bullets or ["No details supplied."], theme, size=16)


def render_presentation(
    spec_path: Path,
    output_path: Path,
    *,
    style: str = "Executive dark",
) -> dict[str, Any]:
    """Render a constrained deck JSON specification to a polished PPTX file."""
    try:
        payload = json.loads(spec_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise PresentationRenderError(f"Could not read deck-content.json: {exc}") from exc
    if not isinstance(payload, dict):
        raise PresentationRenderError("deck-content.json must contain an object.")
    deck_title = _clean_text(payload.get("deck_title"), field="deck_title", maximum=160)
    raw_slides = payload.get("slides")
    if not isinstance(raw_slides, list) or not 2 <= len(raw_slides) <= 20:
        raise PresentationRenderError("A deck must contain between 2 and 20 slides.")
    if not all(isinstance(slide, dict) for slide in raw_slides):
        raise PresentationRenderError("Every slide must be an object.")

    theme = THEMES.get(style, THEMES["Executive dark"])
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    renderers = {
        "title": _add_title_slide,
        "section": _add_section_slide,
        "content": _add_content_slide,
        "metrics": _add_metrics_slide,
        "comparison": _add_comparison_slide,
    }
    for index, slide_data in enumerate(raw_slides, start=1):
        slide_type = _optional_text(slide_data.get("type"), 30).lower() or "content"
        renderer = renderers.get(slide_type)
        if renderer is None:
            raise PresentationRenderError(f"Unsupported slide type: {slide_type}.")
        slide = presentation.slides.add_slide(blank_layout)
        _add_background(slide, theme["background"])
        renderer(slide, slide_data, theme)
        if slide_type != "title":
            _add_slide_chrome(
                slide,
                theme=theme,
                deck_title=deck_title,
                slide_number=index,
                slide_count=len(raw_slides),
                sources=_text_list(slide_data.get("sources"), maximum_items=3, maximum_length=100),
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return {
        "filename": output_path.name,
        "slide_count": len(raw_slides),
        "style": style if style in THEMES else "Executive dark",
        "file_size": output_path.stat().st_size,
    }
