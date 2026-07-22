from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.file_preview import (
    MAX_PREVIEW_FILE_BYTES,
    MAX_TABLE_COLUMNS,
    MAX_TABLE_ROWS,
    build_file_preview,
)


def preview(path: Path, content_type: str):
    return build_file_preview(
        path,
        filename=path.name,
        content_type=content_type,
        file_size=path.stat().st_size,
    )


def test_previews_and_formats_json_without_executing_content(tmp_path: Path) -> None:
    source = tmp_path / "brief.json"
    source.write_text('{"status":"ready","items":[1,2]}', encoding="utf-8")

    result = preview(source, "application/json")

    assert result.kind == "text"
    assert result.text == '{\n  "status": "ready",\n  "items": [\n    1,\n    2\n  ]\n}'
    assert result.truncated is False


def test_csv_preview_is_bounded_by_rows_and_columns(tmp_path: Path) -> None:
    source = tmp_path / "large.csv"
    source.write_text(
        "\n".join(
            ",".join(f"r{row}c{column}" for column in range(MAX_TABLE_COLUMNS + 2))
            for row in range(MAX_TABLE_ROWS + 2)
        ),
        encoding="utf-8",
    )

    result = preview(source, "text/csv")

    assert result.kind == "table"
    assert len(result.tables) == 1
    assert len(result.tables[0].rows) == MAX_TABLE_ROWS
    assert len(result.tables[0].rows[0]) == MAX_TABLE_COLUMNS
    assert result.tables[0].truncated is True
    assert result.truncated is True


def test_xlsx_preview_preserves_sheet_names_and_calculated_values(tmp_path: Path) -> None:
    source = tmp_path / "scorecard.xlsx"
    workbook = Workbook()
    overview = workbook.active
    overview.title = "Overview"
    overview.append(["Metric", "Value"])
    overview.append(["Revenue", 125000])
    detail = workbook.create_sheet("Detail")
    detail.append(["Region", "Status"])
    detail.append(["East", "On track"])
    workbook.save(source)

    result = preview(
        source,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

    assert result.kind == "table"
    assert [table.name for table in result.tables] == ["Overview", "Detail"]
    assert result.tables[0].rows[1] == ["Revenue", "125000"]
    assert result.tables[1].rows[1] == ["East", "On track"]


def test_docx_preview_reads_paragraphs_and_tables(tmp_path: Path) -> None:
    source = tmp_path / "plan.docx"
    document = Document()
    document.add_heading("Launch plan", level=1)
    document.add_paragraph("Review the rollout risks before approval.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Owner"
    table.rows[0].cells[1].text = "Finance"
    document.save(source)

    result = preview(
        source,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    assert result.kind == "text"
    assert "Launch plan" in (result.text or "")
    assert "Review the rollout risks" in (result.text or "")
    assert "Owner\tFinance" in (result.text or "")


def test_pptx_preview_labels_slide_text(tmp_path: Path) -> None:
    source = tmp_path / "briefing.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly update"
    slide.placeholders[1].text = "Growth improved while risk stayed stable."
    presentation.save(source)

    result = preview(
        source,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    assert result.kind == "text"
    assert "Slide 1" in (result.text or "")
    assert "Quarterly update" in (result.text or "")
    assert "Growth improved" in (result.text or "")


def test_images_use_an_authenticated_blob_preview_without_inline_data(tmp_path: Path) -> None:
    source = tmp_path / "diagram.png"
    source.write_bytes(b"not-decoded-by-the-preview-contract")

    result = preview(source, "image/png")

    assert result.kind == "image"
    assert result.text is None
    assert result.tables == []


def test_large_and_legacy_files_fail_closed_with_a_useful_fallback(tmp_path: Path) -> None:
    legacy = tmp_path / "budget.xls"
    legacy.write_bytes(b"legacy")
    large = tmp_path / "archive.txt"
    large.write_text("small on disk", encoding="utf-8")

    legacy_result = preview(legacy, "application/vnd.ms-excel")
    large_result = build_file_preview(
        large,
        filename=large.name,
        content_type="text/plain",
        file_size=MAX_PREVIEW_FILE_BYTES + 1,
    )

    assert legacy_result.kind == "unavailable"
    assert "older Office format" in legacy_result.detail
    assert large_result.kind == "unavailable"
    assert "too large" in large_result.detail
