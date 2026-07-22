import json
import uuid
from pathlib import Path

import pytest
from pptx import Presentation

from app.presentation_renderer import PresentationRenderError, render_presentation
from app.routers import agents as agents_router
from app.work_packs import normalize_work_pack


def _deck_spec() -> dict[str, object]:
    return {
        "deck_title": "Customer retention plan",
        "slides": [
            {
                "type": "title",
                "title": "Retention can return to target in 90 days",
                "subtitle": "A focused plan for the executive team",
            },
            {
                "type": "content",
                "title": "Onboarding friction drives the largest preventable loss",
                "takeaway": "Three early-life interventions address the dominant risk.",
                "bullets": [
                    "Instrument activation milestones",
                    "Add proactive outreach for stalled accounts",
                    "Make the first-value path visible to account teams",
                ],
                "sources": ["retention-analysis.xlsx"],
            },
            {
                "type": "metrics",
                "title": "The recovery plan has measurable leading indicators",
                "metrics": [
                    {"value": "72%", "label": "Activation target", "context": "Within 14 days"},
                    {"value": "<5d", "label": "Time to first value", "context": "Median target"},
                    {"value": "+8pt", "label": "Renewal lift", "context": "Planning assumption"},
                ],
                "sources": ["retention-analysis.xlsx", "customer-notes.docx"],
            },
            {
                "type": "comparison",
                "title": "A focused rollout beats a broad program",
                "columns": [
                    {"heading": "Focused rollout", "bullets": ["Fast signal", "Clear ownership"]},
                    {"heading": "Broad program", "bullets": ["More coverage", "Slower learning"]},
                ],
                "sources": ["rollout-options.pdf"],
            },
        ],
    }


def test_renderer_creates_openable_widescreen_powerpoint(tmp_path: Path) -> None:
    spec_path = tmp_path / "deck-content.json"
    output_path = tmp_path / "briefing-deck.pptx"
    spec_path.write_text(json.dumps(_deck_spec()), encoding="utf-8")

    metadata = render_presentation(spec_path, output_path, style="Warm editorial")

    assert metadata["slide_count"] == 4
    assert metadata["style"] == "Warm editorial"
    assert metadata["file_size"] > 10_000
    presentation = Presentation(output_path)
    assert len(presentation.slides) == 4
    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9, rel=0.01)
    all_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Retention can return to target in 90 days" in all_text
    assert "Sources: retention-analysis.xlsx" in all_text
    assert "Focused rollout" in all_text


def test_renderer_rejects_unsupported_or_unbounded_specs(tmp_path: Path) -> None:
    spec_path = tmp_path / "deck-content.json"
    output_path = tmp_path / "briefing-deck.pptx"
    spec_path.write_text(
        json.dumps(
            {
                "deck_title": "Unsafe deck",
                "slides": [
                    {"type": "title", "title": "Title"},
                    {"type": "embedded-code", "title": "Unsupported"},
                ],
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(PresentationRenderError, match="Unsupported slide type"):
        render_presentation(spec_path, output_path)

    spec_path.write_text(json.dumps({"deck_title": "Too short", "slides": []}), encoding="utf-8")
    with pytest.raises(PresentationRenderError, match="between 2 and 20"):
        render_presentation(spec_path, output_path)


def test_execution_postprocessor_renders_from_persisted_deck_spec(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    user_id = uuid.uuid4()
    execution_id = uuid.uuid4()
    uploads_root = tmp_path / "uploads"
    spec_path = uploads_root / str(user_id) / "generated" / str(execution_id) / "deck-content.json"
    spec_path.parent.mkdir(parents=True)
    spec_path.write_text(json.dumps(_deck_spec()), encoding="utf-8")
    monkeypatch.setattr(agents_router, "PROJECT_ROOT", tmp_path)
    monkeypatch.setattr(agents_router, "UPLOADS_ROOT", uploads_root)
    work_pack = normalize_work_pack(
        "presentation-creation",
        {
            "id": "presentation-creation",
            "answers": {
                "audience": "Executive team",
                "purpose": "Approve the retention recovery plan",
                "slide_count": "5–7 slides",
                "story": ["Recommendation", "Evidence", "Next steps"],
                "style": "Clean light",
                "speaker_notes": True,
            },
        },
    )

    artifact, result = agents_router._render_presentation_work_pack(
        work_pack=work_pack,
        artifact_paths=[str(spec_path.relative_to(tmp_path))],
        user_id=user_id,
        execution_id=execution_id,
    )

    assert artifact is not None
    assert artifact["filename"] == "briefing-deck.pptx"
    assert artifact["content_type"].endswith("presentationml.presentation")
    assert result is not None
    assert result["success"] is True
    assert result["source"] == "workerbee-renderer"
    rendered_path = tmp_path / artifact["storage_path"]
    assert rendered_path.exists()
    assert len(Presentation(rendered_path).slides) == 4
